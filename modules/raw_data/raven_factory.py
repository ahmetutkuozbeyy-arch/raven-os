import os
import json
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation

# ================= AYARLAR =================
# Google AI Studio'dan aldığın API Anahtarını buraya yapıştır
API_KEY = "BURAYA_GEMINI_API_KEY_YAPISTIR" 

# Klasör Yolları
INPUT_FOLDER = "raw_data"       # Dosyaların olduğu klasör
OUTPUT_FOLDER = "modules"       # JSON'ların kaydedileceği klasör
INDEX_FILE = "kutuphane.json"   # Kütüphane listesi

# Yapay Zeka Ayarları
genai.configure(api_key=API_KEY)

generation_config = {
  "temperature": 0.7,
  "top_p": 0.95,
  "top_k": 64,
  "max_output_tokens": 8192,
  "response_mime_type": "application/json",
}

SYSTEM_INSTRUCTION = """
Sen "RA'VEN OS" isimli siber-medikal eğitim simülasyonunun baş içerik mimarısın. 
Görevin, verilen ham tıbbi veriyi (ders notu, slayt metni) analiz edip sistemin işleyebileceği JSON formatına dönüştürmektir.

ÇIKTI SADECE VE SADECE JSON OLMALIDIR.

JSON ŞEMASI:
{
  "meta": { "title": "KONU BAŞLIĞI", "tag": "DERS ADI" },
  "modules": {
    "id": "benzersiz_id_ingilizce", 
    "code": "AUTO-G", 
    "title": "KONU BAŞLIĞI", 
    "author": "RA'VEN AI",
    "tag": "DERS ADI", 
    "class": "3", 
    "block": "GENEL", 
    "url": "modules/dosya_adi.json",
    "summary": { "text": "<div class='summary-container'><div class='sum-card sum-clinic'><div class='sum-title'>TERMİNAL VERİSİ</div><div class='sum-body'>...ÖNEMLİ BİLGİLER...</div></div></div>" },
    "cases": [ 
        { "id": "c1", "title": "GÖREV 1: ...", "desc": "...", "start": "s1", "scenes": { "s1": { "text": "...", "opts": [{"txt":"...", "type":"correct/wrong", "go/action":"..."}] } } } 
    ],
    "quiz": [ { "q": "...", "a": ["..."], "correct": 0 } ],
    "match": [ { "A": "...", "B": "..." } ],
    "decypher": [ { "text": "...", "decoys": ["..."] } ]
  }
}
NOT: id kısmı dosya adıyla uyumlu olmalı.
"""

model = genai.GenerativeModel(
  model_name="gemini-1.5-flash",
  generation_config=generation_config,
  system_instruction=SYSTEM_INSTRUCTION,
)

# --- DOSYA OKUMA FONKSİYONLARI ---
def read_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        print(f"❌ PDF Okuma Hatası: {e}")
        return None

def read_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"❌ PPTX Okuma Hatası: {e}")
        return None

def read_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        print(f"❌ TXT Okuma Hatası: {e}")
        return None

# --- ANA İŞLEMLER ---
def generate_module(text_content, filename):
    print(f"🤖 RA'VEN (Gemini) Analiz Ediyor: {filename}...")
    try:
        # Çok uzun metinleri kırpmak gerekebilir ama Flash modeli genelde 1 kitabı bile alır.
        chat_session = model.start_chat(history=[])
        response = chat_session.send_message(f"DOSYA ADI: {filename}\n\nİÇERİK:\n{text_content}")
        return json.loads(response.text.strip())
    except Exception as e:
        print(f"❌ AI Hatası ({filename}): {e}")
        return None

def update_library_index(new_module_data, json_filename):
    if not os.path.exists(INDEX_FILE):
        library = []
    else:
        with open(INDEX_FILE, 'r', encoding='utf-8') as f:
            try:
                library = json.load(f)
            except:
                library = []

    module_info = new_module_data.get("modules", {})
    if not module_info: return

    entry = {
        "id": module_info.get("id"),
        "code": module_info.get("code", "GEM-01"),
        "title": module_info.get("title"),
        "author": "RA'VEN AI",
        "tag": module_info.get("tag"),
        "class": module_info.get("class", "3"),
        "block": module_info.get("block", "GENEL"),
        "url": f"modules/{json_filename}"
    }

    existing_idx = next((i for i, item in enumerate(library) if item["url"] == entry["url"]), -1)
    
    if existing_idx != -1:
        library[existing_idx] = entry
        print(f"🔄 Kütüphane güncellendi: {entry['title']}")
    else:
        library.append(entry)
        print(f"✅ Kütüphaneye eklendi: {entry['title']}")

    with open(INDEX_FILE, 'w', encoding='utf-8') as f:
        json.dump(library, f, indent=2, ensure_ascii=False)

def main():
    if not os.path.exists(OUTPUT_FOLDER): os.makedirs(OUTPUT_FOLDER)
    if not os.path.exists(INPUT_FOLDER): os.makedirs(INPUT_FOLDER)

    # Desteklenen uzantılar
    files = [f for f in os.listdir(INPUT_FOLDER) if f.endswith(('.txt', '.pdf', '.pptx'))]
    
    if not files:
        print("⚠️ 'raw_data' klasörü boş veya desteklenen dosya yok (txt, pdf, pptx).")
        return

    print(f"🚀 FABRİKA BAŞLATILIYOR... {len(files)} dosya kuyrukta.")

    for file in files:
        file_path = os.path.join(INPUT_FOLDER, file)
        raw_text = ""

        # Dosya tipine göre okuma
        if file.endswith(".pdf"):
            raw_text = read_pdf(file_path)
        elif file.endswith(".pptx"):
            raw_text = read_pptx(file_path)
        elif file.endswith(".txt"):
            raw_text = read_txt(file_path)

        if raw_text and len(raw_text) > 50: # En az 50 karakter veri varsa işle
            module_json = generate_module(raw_text, file)
            
            if module_json:
                output_filename = os.path.splitext(file)[0] + ".json"
                output_path = os.path.join(OUTPUT_FOLDER, output_filename)
                
                with open(output_path, 'w', encoding='utf-8') as f:
                    json.dump(module_json, f, indent=2, ensure_ascii=False)
                
                update_library_index(module_json, output_filename)
        else:
            print(f"⚠️ Dosya boş veya okunamadı: {file}")

    print("\n🏁 TÜM İŞLEMLER TAMAMLANDI.")

if __name__ == "__main__":
    main()